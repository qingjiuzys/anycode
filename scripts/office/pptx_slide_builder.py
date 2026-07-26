"""Build editable native OOXML slides from manifest + brand tokens."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _rgb(tokens: dict, key: str) -> RGBColor:
    from brand_kit import color_from_tokens

    r, g, b = color_from_tokens(tokens, key)
    return RGBColor(r, g, b)


def _box(slide, left, top, width, height, text, *, size=18, bold=False, color_key="body_text", align=PP_ALIGN.LEFT, tokens=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text or ""
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text or ""
    run.font.size = Pt(size)
    run.font.bold = bold
    if tokens is not None:
        run.font.color.rgb = _rgb(tokens, color_key)
    return box


def _bg(slide, prs, tokens, color_key: str):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(tokens, color_key)
    shape.line.fill.background()
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def _accent_bar(slide, prs, tokens):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(tokens, "accent")
    bar.line.fill.background()


def _footer(slide, left_text: str, right_text: str, tokens):
    _box(slide, Inches(0.6), Inches(6.9), Inches(5.5), Inches(0.35), left_text, size=9, color_key="muted_text", tokens=tokens)
    _box(slide, Inches(7.0), Inches(6.9), Inches(5.5), Inches(0.35), right_text, size=9, color_key="muted_text", align=PP_ALIGN.RIGHT, tokens=tokens)


def _add_chart(slide, spec: dict, left=2.0, top=1.5, width=9.0, height=4.5):
    ctype = (spec.get("type") or "bar").lower()
    categories = spec.get("categories") or []
    series_list = spec.get("series") or []
    if not categories or not series_list:
        return
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list[:3]:
        name = s.get("name", "Series")
        values = s.get("data") or s.get("values") or []
        chart_data.add_series(name, tuple(values))
    chart_type = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED)
    slide.shapes.add_chart(
        chart_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    )


def _add_chart_image(slide, png_path: Path, left=2.0, top=1.5, width=9.0, height=4.5):
    if png_path.is_file():
        slide.shapes.add_picture(str(png_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def build_cover(slide, prs, data: dict, tokens: dict):
    _bg(slide, prs, tokens, "primary")
    _accent_bar(slide, prs, tokens)
    _box(slide, Inches(0.9), Inches(1.8), Inches(11.5), Inches(1.4), data.get("title", ""), size=44, bold=True, color_key="header_text", tokens=tokens)
    _box(slide, Inches(0.9), Inches(3.2), Inches(11), Inches(0.9), data.get("subtitle", ""), size=22, color_key="surface_alt", tokens=tokens)
    _box(slide, Inches(0.9), Inches(4.2), Inches(10), Inches(0.5), data.get("meta", ""), size=14, color_key="surface_alt", tokens=tokens)
    chips = data.get("chips") or []
    for i, chip in enumerate(chips[:3]):
        x = Inches(0.9 + i * 4.0)
        _box(slide, x, Inches(5.0), Inches(3.5), Inches(0.45), chip.get("number", ""), size=28, bold=True, color_key="accent", tokens=tokens)
        _box(slide, x, Inches(5.5), Inches(3.5), Inches(0.7), chip.get("text", ""), size=12, color_key="header_text", tokens=tokens)
    _footer(slide, data.get("footer_left", ""), data.get("footer_right", ""), tokens)


def build_section(slide, prs, data: dict, tokens: dict):
    _bg(slide, prs, tokens, "surface_alt")
    _box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.0), data.get("title", ""), size=40, bold=True, color_key="primary", tokens=tokens)
    _box(slide, Inches(0.8), Inches(3.5), Inches(5.2), Inches(2.5), data.get("subtitle", ""), size=20, color_key="body_text", tokens=tokens)
    y = Inches(1.2)
    for i, item in enumerate((data.get("agenda") or [])[:6]):
        _box(slide, Inches(7.0), y + Inches(i * 0.75), Inches(5.5), Inches(0.65), f"{i+1}. {item}", size=20, color_key="primary", tokens=tokens)
    _footer(slide, data.get("footer_left", ""), data.get("footer_right", ""), tokens)


def build_content(slide, prs, data: dict, tokens: dict):
    _bg(slide, prs, tokens, "FFFFFF")
    _accent_bar(slide, prs, tokens)
    _box(slide, Inches(0.7), Inches(0.4), Inches(11.5), Inches(0.8), data.get("title", ""), size=32, bold=True, color_key="primary", tokens=tokens)
    bullets = data.get("bullets") or []
    y = Inches(1.35)
    for i, b in enumerate(bullets[:6]):
        line = b.get("text", "")
        if b.get("stat"):
            line = f"{b['stat']}  {line}"
        _box(slide, Inches(0.85), y + Inches(i * 0.62), Inches(7.0), Inches(0.55), f"• {line}", size=17, color_key="body_text", tokens=tokens)
        if b.get("source"):
            _box(slide, Inches(1.0), y + Inches(i * 0.62) + Inches(0.38), Inches(6.8), Inches(0.3), b["source"], size=11, color_key="muted_text", tokens=tokens)
    if data.get("panel_title") or data.get("panel_body"):
        _box(slide, Inches(8.1), Inches(1.35), Inches(4.5), Inches(0.4), data.get("panel_title", "Summary"), size=14, bold=True, color_key="primary", tokens=tokens)
        _box(slide, Inches(8.1), Inches(1.85), Inches(4.5), Inches(2.0), data.get("panel_body", ""), size=13, color_key="body_text", tokens=tokens)
    if data.get("quote"):
        _box(slide, Inches(8.1), Inches(4.0), Inches(4.5), Inches(2.2), data.get("quote", ""), size=13, color_key="header_text", tokens=tokens)
    _footer(slide, data.get("footer_left", ""), data.get("footer_right", ""), tokens)


def build_metrics(slide, prs, data: dict, tokens: dict, workspace: Path | None = None):
    _bg(slide, prs, tokens, "surface")
    _accent_bar(slide, prs, tokens)
    _box(slide, Inches(0.7), Inches(0.35), Inches(11.5), Inches(0.7), data.get("title", ""), size=30, bold=True, color_key="primary", align=PP_ALIGN.CENTER, tokens=tokens)
    charts = data.get("charts") or []
    if charts:
        spec = charts[0]
        png = spec.get("png") or spec.get("image")
        # Prefer a native editable chart whenever data is present; the PNG is
        # only a fallback for charts that arrive as pure images.
        if spec.get("categories") and spec.get("series"):
            _add_chart(slide, spec)
        elif png and workspace:
            _add_chart_image(slide, workspace / png)
        else:
            _add_chart(slide, spec)
        _footer(slide, data.get("footer_left", ""), data.get("footer_right", ""), tokens)
        return
    stats = data.get("stats") or []
    positions = [(0.7, 1.2), (4.5, 1.2), (8.3, 1.2), (0.7, 3.8), (4.5, 3.8), (8.3, 3.8)]
    for i, (x, y) in enumerate(positions):
        if i >= len(stats):
            break
        s = stats[i]
        _box(slide, Inches(x), Inches(y), Inches(3.5), Inches(0.55), s.get("number", ""), size=34, bold=True, color_key="primary", tokens=tokens)
        _box(slide, Inches(x), Inches(y + 0.55), Inches(3.5), Inches(0.45), s.get("label", ""), size=14, bold=True, color_key="body_text", tokens=tokens)
        _box(slide, Inches(x), Inches(y + 1.0), Inches(3.5), Inches(0.9), s.get("detail", ""), size=12, color_key="body_text", tokens=tokens)
        _box(slide, Inches(x), Inches(y + 1.85), Inches(3.5), Inches(0.35), s.get("source", ""), size=10, color_key="muted_text", tokens=tokens)
    _footer(slide, data.get("footer_left", ""), data.get("footer_right", ""), tokens)


def build_closing(slide, prs, data: dict, tokens: dict):
    _bg(slide, prs, tokens, "primary")
    _accent_bar(slide, prs, tokens)
    _box(slide, Inches(0.9), Inches(0.8), Inches(11), Inches(0.9), data.get("title", "Next Steps"), size=36, bold=True, color_key="header_text", tokens=tokens)
    _box(slide, Inches(0.9), Inches(1.6), Inches(11), Inches(0.5), data.get("subtitle", ""), size=18, color_key="surface_alt", tokens=tokens)
    actions = data.get("actions") or data.get("bullets") or []
    y = Inches(2.3)
    for i, a in enumerate(actions[:5]):
        txt = a.get("text", "") if isinstance(a, dict) else str(a)
        _box(slide, Inches(0.95), y + Inches(i * 0.55), Inches(7.0), Inches(0.5), f"✓ {txt}", size=16, color_key="header_text", tokens=tokens)
    contact = data.get("meta") or data.get("quote") or ""
    _box(slide, Inches(8.2), Inches(2.3), Inches(4.3), Inches(3.5), contact, size=14, color_key="header_text", tokens=tokens)
    _footer(slide, data.get("footer_left", ""), data.get("footer_right", ""), tokens)


BUILDERS = {
    "cover": build_cover,
    "section": build_section,
    "content": build_content,
    "metrics": build_metrics,
    "closing": build_closing,
}


def build_deck(manifest: dict, tokens: dict, workspace: Path | None = None) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    ws = workspace
    for slide_data in manifest.get("slides") or []:
        st = slide_data.get("type", "content")
        slide = prs.slides.add_slide(blank)
        if st == "metrics":
            build_metrics(slide, prs, slide_data, tokens, ws)
        else:
            builder = BUILDERS.get(st, build_content)
            builder(slide, prs, slide_data, tokens)
    return prs
