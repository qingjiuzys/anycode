#!/usr/bin/env python3
"""Render PPTX slide thumbnails — reuse evidence/ PNGs or screenshot sibling slides/*.html."""
from __future__ import annotations

import shutil
import subprocess
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: pip install python-pptx", file=sys.stderr)
    sys.exit(2)

REPO = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(REPO / "scripts" / "office"))


def soffice() -> str | None:
    for name in ("soffice", "libreoffice"):
        p = shutil.which(name)
        if p:
            return p
    return None


def reuse_or_rescreenshot(pptx: Path, out_dir: Path) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    existing = sorted(out_dir.glob("slide-*.png"))
    if existing:
        return existing

    # Prefer re-screenshot from HTML source (pixel-perfect)
    for slides_dir in (pptx.parent / "slides", pptx.parent.parent / "slides"):
        html_files = sorted(slides_dir.glob("*.html")) if slides_dir.is_dir() else []
        if len(html_files) >= 2:
            from screenshot_slide_html import screenshot_html

            paths: list[Path] = []
            for i, hf in enumerate(html_files, start=1):
                dest = out_dir / f"slide-{i:02d}.png"
                screenshot_html(hf, dest)
                paths.append(dest)
            return paths

    sibling_evidence = pptx.parent / "evidence"
    if sibling_evidence.is_dir():
        copied: list[Path] = []
        for i, src in enumerate(sorted(sibling_evidence.glob("slide-*.png")), start=1):
            dest = out_dir / f"slide-{i:02d}.png"
            dest.write_bytes(src.read_bytes())
            copied.append(dest)
        if copied:
            return copied

    return []


def render_soffice(pptx: Path, out_dir: Path) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        soffice() or "",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(out_dir),
        str(pptx),
    ]
    subprocess.run(cmd, check=True, capture_output=True)
    pdf = out_dir / (pptx.stem + ".pdf")
    if not pdf.is_file():
        raise RuntimeError(f"LibreOffice did not produce {pdf}")
    ppm = shutil.which("pdftoppm")
    if ppm:
        subprocess.run(
            [ppm, "-png", str(pdf), str(out_dir / "slide")],
            check=True,
            capture_output=True,
        )
        return sorted(out_dir.glob("slide-*.png"))
    return [pdf]


def render_pil_fallback(pptx: Path, out_dir: Path) -> list[Path]:
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        return []
    out_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(pptx))
    paths: list[Path] = []
    for i, slide in enumerate(prs.slides, start=1):
        img = Image.new("RGB", (1280, 720), color=(238, 243, 248))
        draw = ImageDraw.Draw(img)
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip()[:120])
            elif shape.shape_type == 13:  # picture
                texts.append("[visual slide]")
        y = 40
        draw.rectangle([0, 0, 1280, 80], fill=(27, 58, 92))
        draw.text((30, 24), texts[0] if texts else f"Slide {i}", fill=(255, 255, 255))
        for t in texts[1:6]:
            y += 36
            draw.text((40, y), t[:100], fill=(64, 64, 64))
        dest = out_dir / f"slide-{i:02d}.png"
        img.save(dest)
        paths.append(dest)
    return paths


def main() -> int:
    if len(sys.argv) < 2:
        print("usage: render_pptx_evidence.py deck.pptx [evidence_dir]", file=sys.stderr)
        return 1
    pptx = Path(sys.argv[1]).resolve()
    out_dir = Path(sys.argv[2]).resolve() if len(sys.argv) > 2 else pptx.parent / "evidence"

    paths = reuse_or_rescreenshot(pptx, out_dir)
    if not paths and soffice():
        paths = render_soffice(pptx, out_dir)
    if not paths:
        paths = render_pil_fallback(pptx, out_dir)
    if not paths:
        print("WARN: no thumbnails rendered", file=sys.stderr)
        return 1
    for p in paths:
        print(p)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
