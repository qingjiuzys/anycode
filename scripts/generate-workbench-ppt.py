#!/usr/bin/env python3
"""Generate anyCode Digital Workbench functional guide as PowerPoint."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
GUIDE = ROOT / "docs" / "workbench-functional-guide"
SHOTS = GUIDE / "screenshots"
OUT = GUIDE / "anyCode数字工作台功能说明.pptx"

ACCENT = RGBColor(0x1A, 0x73, 0xE8)
TEXT = RGBColor(0x20, 0x21, 0x24)
MUTED = RGBColor(0x5F, 0x63, 0x68)


def set_title(slide, title: str, subtitle: str | None = None) -> None:
    slide.shapes.title.text = title
    p = slide.shapes.title.text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT
    if subtitle and slide.placeholders[1]:
        body = slide.placeholders[1].text_frame
        body.clear()
        p2 = body.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = MUTED


def add_bullets(slide, items: list[str], left=0.6, top=1.6, width=8.8, height=5.0) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)


def add_image_slide(prs, title: str, image_name: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(layout)
    # title bar
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = ACCENT

    img_path = SHOTS / image_name
    if img_path.is_file():
        slide.shapes.add_picture(
            str(img_path),
            Inches(0.45),
            Inches(1.05),
            width=Inches(6.2),
        )

    add_bullets(slide, bullets, left=6.85, top=1.05, width=2.9, height=5.8)


def main() -> None:
    if not SHOTS.is_dir():
        raise SystemExit(f"screenshots not found: {SHOTS}")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # cover
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_title(
        slide,
        "anyCode 数字工作台",
        "功能说明 · v0.2.0\nhttp://127.0.0.1:43180  ·  anycode dashboard --open",
    )

    # positioning
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, "1. 产品定位", None)
    add_bullets(
        slide,
        [
            "本地 Web 看板：统一管理项目、会话、自动化与产出",
            "AI 任务执行、审批、Gate 门禁与信任度",
            "Agent / Skills / MCP / Browser 连接器配置",
            "数据：~/.anycode/projects.db + SSE 实时事件流",
        ],
        top=1.5,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, "2. 启动与登录", None)
    add_bullets(
        slide,
        [
            "浏览器：anycode dashboard --open",
            "macOS：anyCode.app 内置启动工作台",
            "账户：local@anycode（本地 trusted 免密）",
            "右上角：中/英、深浅色；SSE 已连接 = 实时流正常",
        ],
        top=1.5,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, "3. 侧栏导航一览", None)
    add_bullets(
        slide,
        [
            "总览 — 运维摘要、快捷对话、事件流",
            "项目 — registry、信任度、扫描/新建",
            "会话 — Web 聊天、追问、产物",
            "自动化 — Cron、项目护栏",
            "资产 — 文件写入跟踪与验证",
            "报告 / 审计 — 导出与配置变更",
            "Agent·Skills — Starter 包、治理",
            "设置 — 模型、MCP、Browser、Doctor",
        ],
        top=1.35,
        height=5.5,
    )

    slides_with_shots = [
        (
            "4. 总览首页",
            "01-home.png",
            [
                "快捷对话：选项目即开新会话",
                "运维摘要：阻断/审批/预算",
                "活跃项目 + 实时事件流",
                "事件标题 UI 层中文化",
            ],
        ),
        (
            "5. 项目列表",
            "02-projects.png",
            [
                "扫描 / 新建 / 归档",
                "Flutter 等模板脚手架",
                "信任度进度条",
                "根目录缺失标红",
            ],
        ),
        (
            "6. 项目详情",
            "07-project-detail.png",
            [
                "配置：知识库、Gate、Pipeline",
                "重建索引 / 索引资产",
                "就绪分与 Gate 通过率",
                "生成项目报告",
            ],
        ),
        (
            "7. 会话与 Web 聊天",
            "03-conversations.png",
            [
                "按项目/状态筛选",
                "Agent / 模型 / @Skills",
                "图片附件（Vision 模型）",
                "工具调用与审批可视化",
            ],
        ),
        (
            "8. 自动化",
            "04-automations.png",
            [
                "自然语言 → cron 表达式",
                "写入 orchestration.json",
                "项目护栏：阻断/通知/报告",
                "Cron 运行记录与重试",
            ],
        ),
        (
            "9. 资产",
            "05-assets.png",
            [
                "FileWrite / Edit 跟踪",
                "未验证 / 阻断产物筛选",
                "导出 CSV、查看 hash",
            ],
        ),
        (
            "10. Agent / Skills",
            "06-agents-skills.png",
            [
                "Agent 统计与模型路由",
                "Skills 扫描与治理",
                "Starter 包一键安装",
                "中文场景：日报/周报/纪要",
            ],
        ),
        (
            "11. 设置：MCP 与浏览器",
            "08-settings-notify-mcp-browser.png",
            [
                "MCP 服务器 UI 编辑 config",
                "Browser 连接器（Desktop）",
                "Playwright MCP 无头 Chromium",
                "通知策略与渠道集成",
            ],
        ),
        (
            "12. 设置：偏好与提示词",
            "09-settings-prefs-prompt.png",
            [
                "UI 密度与报告偏好",
                "资产读取策略",
                "System Prompt 预览",
                "CLAUDE.md + Skills 合并视图",
            ],
        ),
    ]

    for title, img, bullets in slides_with_shots:
        add_image_slide(prs, title, img, bullets)

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, "13. 与 CLI 的关系", None)
    add_bullets(
        slide,
        [
            "终端 anycode → 记录事件 → Dashboard 展示",
            "Dashboard Web 聊天 → 同一 AgentRuntime",
            "配置中心：~/.anycode/config.json",
            "详见 docs/run-flow.md",
        ],
        top=1.5,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, "14. 常见问题", None)
    add_bullets(
        slide,
        [
            "空白页：确认 dashboard 进程、刷新浏览器",
            "无数据：终端跑一次任务或「扫描项目」",
            "知识库向量：Desktop 自带；CLI 需 knowledge-embeddings",
            "Browser 不可用：设置中启用连接器并重启",
        ],
        top=1.5,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_title(slide, "谢谢", "anyCode · Digital Workbench\n文档：docs/workbench-functional-guide/")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
