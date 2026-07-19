#!/usr/bin/env python3
"""Verify 08-complex-delivery v2: semantic multi-artifact office + Rust + git + manifest."""

from __future__ import annotations

import json
import hashlib
import re
import subprocess
import sys
from html import unescape
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from office_verify import (  # noqa: E402
    Finding,
    ScoreResult,
    grade_from,
    load_csv_truth,
    verify_docx,
    verify_pptx,
    verify_xlsx,
)

HARNESS = Path(__file__).resolve().parents[2]
V2 = json.loads((HARNESS / "shared/quality/complex_v2.json").read_text(encoding="utf-8"))
MIN_SLIDES = int(V2["min_ppt_slides"])
MIN_PPT_THEMES = int(V2.get("min_ppt_themes", 8))
TRUTH = V2["truth"]
PPT_THEMES = V2.get("ppt_themes", {})
DOCX_SECTIONS = V2.get("docx_sections", [])
QA_TOKENS = V2.get("qa_required_tokens", [])


def run(cmd: list[str], cwd: Path) -> subprocess.CompletedProcess[str]:
    return subprocess.run(cmd, cwd=cwd, capture_output=True, text=True, timeout=300)


def clamp_score(raw: int) -> int:
    return max(0, min(100, raw))


def text_has_number(text: str, n: int | float) -> bool:
    normalized = text.replace(",", "").replace("，", "")
    raw = str(n)
    if raw in normalized:
        return True
    if isinstance(n, float):
        return f"{n:.2f}" in text or f"{n:.1f}" in text
    return False


def strip_html_text(html: str) -> str:
    text = re.sub(r"<script[^>]*>.*?</script>", " ", html, flags=re.I | re.S)
    text = re.sub(r"<style[^>]*>.*?</style>", " ", text, flags=re.I | re.S)
    text = re.sub(r"<[^>]+>", " ", text)
    return unescape(re.sub(r"\s+", " ", text))


def theme_hit(blob: str, keywords: list[str]) -> bool:
    low = blob.lower()
    return any(k.lower() in low for k in keywords)


def extract_ppt_blob(path: Path) -> tuple[list[str], str]:
    from pptx import Presentation

    prs = Presentation(str(path))
    titles: list[str] = []
    chunks: list[str] = []
    for slide in prs.slides:
        slide_bits: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_bits.append(shape.text.strip())
        if slide_bits:
            titles.append(slide_bits[0][:80])
            chunks.extend(slide_bits)
    return titles, "\n".join(chunks)


def verify_ppt_structure(path: Path) -> tuple[int, list[Finding], dict]:
    findings: list[Finding] = []
    if not path.is_file():
        return 0, [Finding("P0", "missing_pptx", f"缺少 {path.name}", 12)], {}
    titles, blob = extract_ppt_blob(path)
    from pptx import Presentation

    n = len(Presentation(str(path)).slides)
    evidence = {"slide_titles": titles, "slides": n}
    score = 100
    if n < MIN_SLIDES:
        findings.append(Finding("P0", "pptx_slides", f"PPT 需 >={MIN_SLIDES} 页，当前 {n}", 15))
        score -= 20
    matched = []
    for theme, keywords in PPT_THEMES.items():
        if theme_hit(blob, keywords):
            matched.append(theme)
    evidence["ppt_themes_matched"] = matched
    if len(matched) < MIN_PPT_THEMES:
        findings.append(
            Finding(
                "P1",
                "pptx_themes",
                f"PPT 主题覆盖不足：{len(matched)}/{MIN_PPT_THEMES}（缺 {set(PPT_THEMES)-set(matched)}）",
                10,
            )
        )
        score -= 12
    return clamp_score(score), findings, evidence


def verify_docx_structure(path: Path, fixtures: Path) -> tuple[int, list[Finding], dict]:
    findings: list[Finding] = []
    if not path.is_file():
        return 0, [Finding("P0", "missing_docx", f"缺少 {path.name}", 12)], {}
    base = verify_docx(path, fixtures)
    from docx import Document

    doc = Document(str(path))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    for t in doc.tables:
        for row in t.rows:
            text += "\n" + " ".join(c.text for c in row.cells)
    matched = [s for s in DOCX_SECTIONS if s in text]
    evidence = {"docx_sections_matched": matched, "base_score": base.score}
    score = base.score
    if len(matched) < 4:
        findings.append(
            Finding("P1", "docx_sections", f"Word 章节不足：{matched}（需背景/发现/影响/根因/行动/负责人）", 10)
        )
        score -= 10
    if not text_has_number(text, TRUTH["total_sales"]):
        findings.append(Finding("P1", "docx_total_sales", "Word 未显式写总销售额 909900", 8))
        score -= 8
    if not text_has_number(text, TRUTH["north_refund_rate_pct"]):
        findings.append(Finding("P1", "docx_north_rate", "Word 未写华北退款率 6.65%", 8))
        score -= 8
    if "负责人" not in text and "owner" not in text.lower():
        findings.append(Finding("P1", "docx_owner", "Word 缺负责人/owner", 6))
        score -= 6
    findings.extend(
        Finding(f.severity, f.code, f.message, f.points)
        for f in base.findings
        if f.severity == "P0"
    )
    return clamp_score(score), findings, evidence


def verify_dashboard_html(path: Path) -> tuple[int, list[Finding], dict]:
    findings: list[Finding] = []
    if not path.is_file():
        return 0, [Finding("P0", "missing_html", f"缺少 {path.name}", 15)], {}
    html = path.read_text(encoding="utf-8", errors="replace")
    plain = strip_html_text(html)
    score = 100
    evidence = {"plain_len": len(plain)}
    if "<!DOCTYPE" not in html.upper() and "<html" not in html.lower():
        findings.append(Finding("P1", "html_doctype", "缺少 HTML5 结构", 8))
        score -= 10
    if not text_has_number(plain, TRUTH["total_sales"]):
        findings.append(Finding("P1", "html_total_sales", "HTML 未含总销售额 909900", 10))
        score -= 12
    if "华北" not in plain or "退款" not in plain:
        findings.append(Finding("P1", "html_north", "HTML 未同时体现华北与退款", 8))
        score -= 10
    has_kpi = bool(re.search(r"kpi|metric|card|指标", html, re.I)) or "909" in plain
    if not has_kpi:
        findings.append(Finding("P2", "html_kpi", "HTML 缺 KPI/卡片结构", 5))
        score -= 5
    has_risk = bool(re.search(r"risk|alert|warning|异常|风险", html + plain, re.I))
    if not has_risk and "华北" in plain:
        evidence["north_in_plain"] = True
    elif not has_risk:
        findings.append(Finding("P2", "html_risk_block", "HTML 缺华北风险区块", 5))
        score -= 5
    marker = "E2E_EXEC_DASHBOARD_MARKER"
    evidence["has_marker"] = marker in html
    if marker not in html:
        findings.append(Finding("P2", "html_marker", f"缺少隐藏 marker {marker}", 5))
        score -= 5
    elif re.search(rf">{marker}<", html) and "display:none" not in html and "display: none" not in html:
        findings.append(Finding("P2", "html_marker_visible", "marker 不应作为可见 CTA", 5))
        score -= 5
    return clamp_score(score), findings, evidence


def verify_qa_signoff(path: Path) -> tuple[int, list[Finding], dict]:
    findings: list[Finding] = []
    if not path.is_file():
        return 0, [Finding("P0", "missing_qa", f"缺少 {path.name}", 12)], {}
    text = path.read_text(encoding="utf-8")
    score = 100
    missing = []
    for token in QA_TOKENS:
        if token.lower() not in text.lower():
            missing.append(token)
            findings.append(Finding("P1", f"qa_missing_{token}", f"qa_signoff 缺 {token}", 5))
            score -= 5
    evidence = {"qa_missing": missing}
    if "- [ ]" not in text and "- [x]" not in text and "- [X]" not in text:
        findings.append(Finding("P2", "qa_checklist", "qa_signoff 应含 checkbox 清单", 4))
        score -= 4
    return clamp_score(score), findings, evidence


def verify_scorecard_xlsx(path: Path, fixtures: Path) -> tuple[int, list[Finding], dict]:
    findings: list[Finding] = []
    if not path.is_file():
        return 0, [Finding("P0", "missing_scorecard", "缺少 regional_scorecard.xlsx", 12)], {}
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    score = 100
    evidence = {"sheets": wb.sheetnames}
    if len(wb.sheetnames) < 3:
        findings.append(Finding("P1", "scorecard_sheets", f"scorecard 需 >=3 工作表，当前 {len(wb.sheetnames)}", 10))
        score -= 15
    truth = load_csv_truth(fixtures)
    blob = " ".join(
        str(c)
        for sn in wb.sheetnames
        for row in wb[sn].iter_rows(values_only=True)
        for c in row
        if c is not None
    )
    for region in ("华东", "华南", "华北"):
        if region not in blob:
            findings.append(Finding("P1", f"scorecard_{region}", f"scorecard 缺区域 {region}", 6))
            score -= 6
    if str(truth["total_sales"]) not in blob.replace(",", ""):
        findings.append(Finding("P1", "scorecard_total", "scorecard 未含总销售额", 8))
        score -= 10
    wb.close()
    return clamp_score(score), findings, evidence


def manifest_has_repo_key(manifest: dict) -> bool:
    for k in V2.get("manifest_repo_aliases", ["code_repo", "git_repo"]):
        if k in manifest and manifest[k]:
            return True
    return False


def cross_check_manifest(manifest: dict, html_plain: str, ppt_blob: str) -> list[Finding]:
    findings: list[Finding] = []
    ts = manifest.get("total_sales")
    if ts != TRUTH["total_sales"]:
        findings.append(Finding("P0", "manifest_total", f"manifest total_sales 应为 {TRUTH['total_sales']}", 12))
    sc = manifest.get("slide_count")
    if isinstance(sc, int) and sc < MIN_SLIDES:
        findings.append(Finding("P1", "manifest_slides", f"manifest slide_count < {MIN_SLIDES}", 6))
    if html_plain and str(TRUTH["total_sales"]) not in html_plain.replace(",", ""):
        if ts == TRUTH["total_sales"]:
            findings.append(Finding("P2", "cross_html_total", "manifest 与 HTML 总销售额不一致", 4))
    if ppt_blob and str(TRUTH["total_sales"]) not in ppt_blob.replace(",", ""):
        findings.append(Finding("P2", "cross_ppt_total", "manifest 与 PPT 总销售额交叉校验失败", 4))
    return findings


def verify_manifest_runtime_truth(
    manifest: dict,
    *,
    cargo_passed: bool | None,
    live_git_head: str | None,
    run_id: str | None,
    profile: str | None,
) -> list[Finding]:
    findings: list[Finding] = []
    if run_id is not None and manifest.get("eval_run_id") != run_id:
        findings.append(Finding("P0", "manifest_run_id", "manifest eval_run_id 不属于本次运行", 20))
    if profile is not None and manifest.get("eval_profile") != profile:
        findings.append(Finding("P0", "manifest_profile", "manifest eval_profile 不属于本次模型", 20))
    if not isinstance(manifest.get("cargo_tests_passed"), bool) or manifest.get("cargo_tests_passed") != cargo_passed:
        findings.append(Finding("P0", "manifest_cargo_truth", "manifest cargo_tests_passed 与实时 cargo test 不一致", 25))
    manifest_head = str(manifest.get("git_head") or "").strip()
    head_matches = bool(
        live_git_head
        and len(manifest_head) >= 7
        and live_git_head.startswith(manifest_head)
    )
    if not head_matches:
        findings.append(Finding("P0", "manifest_git_truth", "manifest git_head 与实时 HEAD 不一致", 20))
    return findings


def verify_ownership(
    workspace: Path,
    owner_file: Path | None,
    run_id: str | None,
    profile: str | None,
    session_id: str | None,
) -> tuple[list[Finding], dict]:
    findings: list[Finding] = []
    evidence: dict = {}
    if owner_file is None or not owner_file.is_file():
        return [Finding("P0", "owner_missing", "缺少评测产物归属记录", 20)], evidence
    try:
        owner = json.loads(owner_file.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError) as exc:
        return [Finding("P0", "owner_invalid", f"产物归属记录无效: {exc}", 20)], evidence

    evidence["owner"] = owner
    expected = {
        "evalRunId": run_id,
        "modelProfile": profile,
        "sessionId": session_id,
        "workspace": str(workspace),
    }
    for key, value in expected.items():
        if value is not None and owner.get(key) != value:
            findings.append(Finding("P0", f"owner_{key}", f"归属 {key} 与本次运行不一致", 20))

    owned = {item.get("path"): item for item in owner.get("artifacts", []) if isinstance(item, dict)}
    for relative in V2["required_artifacts"]:
        path = workspace / relative
        item = owned.get(relative)
        if not path.is_file():
            continue
        if not item:
            findings.append(Finding("P0", "artifact_unowned", f"{relative} 无本次运行归属记录", 20))
            continue
        digest = hashlib.sha256(path.read_bytes()).hexdigest()
        if item.get("sha256") != digest or item.get("fresh") is not True:
            findings.append(Finding("P0", "artifact_owner_mismatch", f"{relative} 非本次运行新产物", 20))
    return findings, evidence


def verify_complex(
    workspace: Path,
    *,
    owner_file: Path | None = None,
    run_id: str | None = None,
    profile: str | None = None,
    session_id: str | None = None,
) -> ScoreResult:
    findings: list[Finding] = []
    dims = {
        "data_accuracy": 30,
        "task_completeness": 30,
        "business_expression": 20,
        "presentation": 15,
        "process_trust": 10,
    }
    evidence: dict = {}
    owner_findings, owner_evidence = verify_ownership(workspace, owner_file, run_id, profile, session_id)
    findings.extend(owner_findings)
    evidence.update(owner_evidence)
    fixtures = workspace / "fixtures"
    artifacts = workspace / "artifacts"
    pack = artifacts / "executive_pack"
    repo = fixtures / "e2e-complex-repo"
    manifest_path = artifacts / "DELIVERY_MANIFEST.json"

    xlsx_ops = pack / "june_ops.xlsx"
    xlsx_score = pack / "regional_scorecard.xlsx"
    docx = pack / "incident_brief.docx"
    pptx = pack / "board_deck.pptx"
    qa_md = pack / "qa_signoff.md"
    html = artifacts / "executive_dashboard.html"

    if xlsx_ops.is_file():
        try:
            sub = verify_xlsx(xlsx_ops, fixtures)
            evidence["xlsx_ops"] = {"score": sub.score, "grade": sub.grade}
            if sub.score < 85:
                findings.append(Finding("P1", "xlsx_ops_quality", f"june_ops 质量分 {sub.score}", 6))
                dims["presentation"] -= 3
        except SystemExit as e:
            findings.append(Finding("P1", "xlsx_ops_verify", str(e), 8))
            dims["data_accuracy"] -= 4
    else:
        findings.append(Finding("P0", "missing_xlsx_ops", "缺少 june_ops.xlsx", 10))
        dims["task_completeness"] -= 10

    sc_score, sc_findings, sc_ev = verify_scorecard_xlsx(xlsx_score, fixtures)
    evidence["scorecard"] = {"score": sc_score, **sc_ev}
    findings.extend(sc_findings)
    if sc_score < 70:
        dims["task_completeness"] -= 8

    doc_score, doc_findings, doc_ev = verify_docx_structure(docx, fixtures)
    evidence["docx"] = {"score": doc_score, **doc_ev}
    findings.extend(doc_findings)
    if doc_score < 85:
        dims["presentation"] -= 3

    ppt_score, ppt_findings, ppt_ev = verify_ppt_structure(pptx)
    evidence["pptx"] = {"score": ppt_score, **ppt_ev}
    evidence["pptx_slides"] = ppt_ev.get("slides")
    findings.extend(ppt_findings)
    if ppt_score < 85:
        dims["presentation"] -= 3

    qa_score, qa_findings, qa_ev = verify_qa_signoff(qa_md)
    evidence["qa_signoff"] = {"score": qa_score, **qa_ev}
    findings.extend(qa_findings)
    if qa_score < 85:
        dims["process_trust"] -= 5

    html_score, html_findings, html_ev = verify_dashboard_html(html)
    evidence["dashboard_html"] = {"score": html_score, **html_ev}
    findings.extend(html_findings)
    html_plain = strip_html_text(html.read_text(encoding="utf-8", errors="replace")) if html.is_file() else ""
    if html_score < 85:
        dims["task_completeness"] -= 5

    _, ppt_blob = extract_ppt_blob(pptx) if pptx.is_file() else ([], "")

    cargo_passed: bool | None = None
    live_git_head: str | None = None
    if not repo.is_dir():
        findings.append(Finding("P0", "missing_repo", "e2e-complex-repo 不存在", 20))
        dims["task_completeness"] -= 20
    else:
        cargo = run(["cargo", "test", "--workspace"], repo)
        evidence["cargo_exit"] = cargo.returncode
        cargo_passed = cargo.returncode == 0
        evidence["cargo_tail"] = (cargo.stdout + cargo.stderr)[-2000:]
        if cargo.returncode != 0:
            findings.append(Finding("P0", "cargo_test_fail", "cargo test --workspace 未通过", 25))
            dims["data_accuracy"] -= 25

        changelog = repo / "CHANGELOG.md"
        evidence["has_changelog"] = changelog.is_file()
        if V2.get("require_changelog") and not changelog.is_file():
            findings.append(Finding("P1", "no_changelog", "仓库缺少 CHANGELOG.md", 8))
            dims["process_trust"] -= 6

        if (repo / ".git").is_dir():
            log = run(["git", "log", "--oneline"], repo)
            head = run(["git", "rev-parse", "HEAD"], repo)
            lines = [ln for ln in log.stdout.strip().splitlines() if ln.strip()]
            evidence["git_commits"] = lines[:5]
            live_git_head = head.stdout.strip() if head.returncode == 0 else None
            evidence["git_head"] = live_git_head or ""

    if manifest_path.is_file():
        try:
            manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
            evidence["manifest"] = manifest
            if not manifest_has_repo_key(manifest):
                findings.append(Finding("P1", "manifest_incomplete", "manifest 缺 code_repo/git_repo", 4))
                dims["process_trust"] -= 3
            for key in V2["manifest_required_keys"]:
                if key == "code_repo":
                    continue
                if key not in manifest:
                    findings.append(Finding("P1", "manifest_incomplete", f"manifest 缺 {key}", 4))
                    dims["process_trust"] -= 2
            findings.extend(
                verify_manifest_runtime_truth(
                    manifest,
                    cargo_passed=cargo_passed,
                    live_git_head=live_git_head,
                    run_id=run_id,
                    profile=profile,
                )
            )
            findings.extend(cross_check_manifest(manifest, html_plain, ppt_blob))
        except json.JSONDecodeError:
            findings.append(Finding("P1", "manifest_invalid", "DELIVERY_MANIFEST.json 无效", 10))
            dims["process_trust"] -= 8
    else:
        findings.append(Finding("P0", "manifest_missing", "缺少 DELIVERY_MANIFEST.json", 12))
        dims["task_completeness"] -= 12

    evidence["complex_v2"] = {"version": V2.get("version"), "min_ppt_slides": MIN_SLIDES}
    raw_score = sum(max(0, v) for v in dims.values())
    score = clamp_score(raw_score)
    grade, pass_gate = grade_from(score, findings)
    return ScoreResult("complex_v2", str(workspace), score, grade, pass_gate, dims, findings, evidence)


def main() -> None:
    import argparse

    p = argparse.ArgumentParser()
    p.add_argument("--workspace", required=True)
    p.add_argument("--owner-file", type=Path)
    p.add_argument("--run-id")
    p.add_argument("--profile")
    p.add_argument("--session-id")
    p.add_argument("--json", action="store_true")
    args = p.parse_args()
    result = verify_complex(
        Path(args.workspace),
        owner_file=args.owner_file,
        run_id=args.run_id,
        profile=args.profile,
        session_id=args.session_id,
    )
    if args.json:
        print(json.dumps(result.to_dict(), ensure_ascii=False, indent=2))
    else:
        print(f"complex_v2 score={result.score} grade={result.grade}")
    raise SystemExit(0 if result.pass_gate else 1)


if __name__ == "__main__":
    main()
