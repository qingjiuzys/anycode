#!/usr/bin/env python3
"""Quality verifier for e2e delivery chain artifacts — scores 0-100 with findings."""

from __future__ import annotations

import csv
import json
import re
import sys
from collections import defaultdict
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any


@dataclass
class Finding:
    severity: str  # P0|P1|P2|P3
    code: str
    message: str
    points: int = 0

    def to_dict(self) -> dict[str, Any]:
        return {"severity": self.severity, "code": self.code, "message": self.message, "points": self.points}


@dataclass
class ScoreResult:
    kind: str
    path: str
    score: int
    grade: str
    pass_gate: bool
    dimensions: dict[str, int] = field(default_factory=dict)
    findings: list[Finding] = field(default_factory=list)
    evidence: dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> dict[str, Any]:
        return {
            "kind": self.kind,
            "artifact": self.path,
            "score": self.score,
            "grade": self.grade,
            "pass": self.pass_gate,
            "dimensions": self.dimensions,
            "findings": [f.to_dict() for f in self.findings],
            "evidence": self.evidence,
        }


def grade_from(score: int, findings: list[Finding]) -> tuple[str, bool]:
    p0 = sum(1 for f in findings if f.severity == "P0")
    p1 = sum(1 for f in findings if f.severity == "P1")
    if p0 > 0 or score < 70:
        return "FAIL", False
    if p1 > 0 or score < 85:
        return "WARN", score >= 70
    return "PASS", True


def load_csv_truth(fixtures_dir: Path) -> dict[str, Any]:
    p = fixtures_dir / "sales_june.csv"
    rows = list(csv.DictReader(p.open(encoding="utf-8")))
    by_date: dict[str, dict[str, int]] = defaultdict(lambda: {"sales": 0, "orders": 0, "refunds": 0})
    by_region: dict[str, dict[str, int]] = defaultdict(lambda: {"sales": 0, "orders": 0, "refunds": 0})
    for r in rows:
        d = r["date"]
        reg = r["region"]
        by_date[d]["sales"] += int(r["sales"])
        by_date[d]["orders"] += int(r["orders"])
        by_date[d]["refunds"] += int(r["refunds"])
        by_region[reg]["sales"] += int(r["sales"])
        by_region[reg]["orders"] += int(r["orders"])
        by_region[reg]["refunds"] += int(r["refunds"])
    total_sales = sum(int(r["sales"]) for r in rows)
    total_orders = sum(int(r["orders"]) for r in rows)
    total_refunds = sum(int(r["refunds"]) for r in rows)
    north = by_region.get("华北", {})
    north_rr = (north.get("refunds", 0) / north.get("orders", 1)) * 100
    return {
        "rows": rows,
        "by_date": dict(by_date),
        "by_region": dict(by_region),
        "total_sales": total_sales,
        "total_orders": total_orders,
        "total_refunds": total_refunds,
        "north_refund_rate": round(north_rr, 2),
    }


def load_md_truth(fixtures_dir: Path) -> str:
    return (fixtures_dir / "report_table.md").read_text(encoding="utf-8")


def digits_in_text(text: str) -> set[str]:
    return set(re.findall(r"\d[\d,]*", text.replace(",", "")))


def verify_xlsx(path: Path, fixtures_dir: Path) -> ScoreResult:
    import openpyxl

    truth = load_csv_truth(fixtures_dir)
    findings: list[Finding] = []
    dims = {"data_accuracy": 30, "task_completeness": 25, "business_expression": 20, "presentation": 15, "process_trust": 10}

    wb = openpyxl.load_workbook(path, read_only=True)
    sheet_names = wb.sheetnames
    evidence: dict[str, Any] = {"sheets": sheet_names}

    # task completeness: two sheets / summary
    has_summary_sheet = any("summary" in n.lower() or "汇总" in n for n in sheet_names)
    has_raw = len(sheet_names) >= 1
    if not has_summary_sheet:
        findings.append(Finding("P1", "xlsx_no_summary_sheet", "缺少 Executive Summary / 汇总工作表", 15))
        dims["task_completeness"] -= 15
    else:
        evidence["summary_sheet"] = next(n for n in sheet_names if "summary" in n.lower() or "汇总" in n)

    ws = wb[wb.sheetnames[0]]
    rows = list(ws.iter_rows(values_only=True))
    evidence["primary_rows"] = len(rows)
    if len(rows) < 4:
        findings.append(Finding("P0", "xlsx_too_few_rows", f"明细行数不足: {len(rows)}", 20))
        dims["data_accuracy"] -= 20

    flat = " ".join(str(c) for r in rows for c in r if c is not None)
    if "华东" not in flat:
        findings.append(Finding("P0", "xlsx_missing_region", "缺少华东区域数据", 15))
        dims["data_accuracy"] -= 15

    # Check summary sheet content if present
    summary_text = ""
    if has_summary_sheet:
        sws = wb[evidence["summary_sheet"]]
        srows = list(sws.iter_rows(values_only=True))
        summary_text = " ".join(str(c) for r in srows for c in r if c is not None)
        evidence["summary_rows"] = len(srows)
        if str(truth["by_region"]["华东"]["sales"]) not in summary_text.replace(",", ""):
            findings.append(Finding("P1", "xlsx_summary_numbers", "汇总表未含华东销售额 393800", 10))
            dims["data_accuracy"] -= 10
        if "华北" not in summary_text and "退款" not in summary_text:
            findings.append(Finding("P1", "xlsx_north_risk", "汇总未标注华北退款风险", 8))
            dims["business_expression"] -= 8
    else:
        dims["business_expression"] -= 10
        findings.append(Finding("P2", "xlsx_raw_only", "仅为 raw dump，无业务汇总表达", 10))

    if str(truth["total_sales"]) not in flat.replace(",", "") and str(truth["total_sales"]) not in summary_text.replace(",", ""):
        findings.append(Finding("P1", "xlsx_total_sales", "未体现总销售额 909900", 8))
        dims["data_accuracy"] -= 8

    score = max(0, sum(max(0, v) for v in dims.values()))
    grade, pass_gate = grade_from(score, findings)
    return ScoreResult("xlsx", str(path), score, grade, pass_gate, dims, findings, evidence)


def verify_docx(path: Path, fixtures_dir: Path) -> ScoreResult:
    from docx import Document

    truth = load_csv_truth(fixtures_dir)
    findings: list[Finding] = []
    dims = {"data_accuracy": 30, "task_completeness": 25, "business_expression": 20, "presentation": 15, "process_trust": 10}

    doc = Document(str(path))
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    text = "\n".join(paras)
    for t in doc.tables:
        for row in t.rows:
            text += "\n" + " ".join(c.text for c in row.cells)

    evidence = {"paragraphs": len(paras), "tables": len(doc.tables)}

    if len(paras) < 3:
        findings.append(Finding("P1", "docx_structure", "段落过少，缺业务结构", 10))
        dims["task_completeness"] -= 10

    if not any(k in text for k in ("结论", "摘要", "总结")):
        findings.append(Finding("P1", "docx_no_conclusion", "缺少结论/摘要段", 12))
        dims["business_expression"] -= 12

    if not any(k in text for k in ("建议", "行动", "下一步")):
        findings.append(Finding("P2", "docx_no_recommendation", "缺少行动建议", 6))
        dims["business_expression"] -= 6

    north_risk = any(k in text for k in ("华北", "退款率", "异常")) and ("华北" in text)
    if not north_risk:
        findings.append(Finding("P1", "docx_north_risk", "未点出华北退款率异常/偏高", 12))
        dims["business_expression"] -= 12

    if "909900" not in text.replace(",", "") and "909,900" not in text:
        findings.append(Finding("P2", "docx_total_sales", "未明确总销售额", 5))
        dims["data_accuracy"] -= 5

    if str(truth["by_region"]["华东"]["sales"]) not in text.replace(",", ""):
        findings.append(Finding("P1", "docx_region_data", "华东销售额与源数据不一致或缺失", 10))
        dims["data_accuracy"] -= 10

    if len(doc.tables) < 1:
        findings.append(Finding("P1", "docx_no_table", "缺少数据表", 10))
        dims["presentation"] -= 10

    score = max(0, sum(max(0, v) for v in dims.values()))
    grade, pass_gate = grade_from(score, findings)
    return ScoreResult("docx", str(path), score, grade, pass_gate, dims, findings, evidence)


def verify_pptx(path: Path, fixtures_dir: Path) -> ScoreResult:
    from pptx import Presentation

    truth = load_csv_truth(fixtures_dir)
    findings: list[Finding] = []
    dims = {"data_accuracy": 30, "task_completeness": 25, "business_expression": 20, "presentation": 15, "process_trust": 10}

    prs = Presentation(str(path))
    n = len(prs.slides)
    all_text = []
    slide_texts = []
    for s in prs.slides:
        st = []
        for sh in s.shapes:
            if hasattr(sh, "text") and sh.text.strip():
                st.append(sh.text.strip())
        slide_texts.append(st)
        all_text.extend(st)
    text = "\n".join(all_text)
    evidence = {"slides": n, "slide_texts": slide_texts}

    if n < 3:
        findings.append(Finding("P0", "pptx_slides", f"页数不足: {n}", 20))
        dims["task_completeness"] -= 20

    if "909900" not in text.replace(",", "") and "909,900" not in text:
        findings.append(Finding("P1", "pptx_total", "未体现总销售额", 8))
        dims["data_accuracy"] -= 8

    if not ("华北" in text and ("退款" in text or "异常" in text)):
        findings.append(Finding("P1", "pptx_north_risk", "结论页缺华北退款风险", 10))
        dims["business_expression"] -= 10

    if not any(k in text.lower() for k in ("sales_june", "fixtures", "数据来源", "csv")):
        findings.append(Finding("P2", "pptx_source", "未注明数据来源", 5))
        dims["presentation"] -= 5

    # detect duplicate KPI blocks
    if text.count("909") > 3 or text.count("总销售") > 2:
        findings.append(Finding("P2", "pptx_repeat_kpi", "多页重复堆叠相同 KPI", 5))
        dims["presentation"] -= 5

    score = max(0, sum(max(0, v) for v in dims.values()))
    grade, pass_gate = grade_from(score, findings)
    return ScoreResult("pptx", str(path), score, grade, pass_gate, dims, findings, evidence)


def verify_html(path: Path, fixtures_dir: Path) -> ScoreResult:
    findings: list[Finding] = []
    dims = {"data_accuracy": 30, "task_completeness": 25, "business_expression": 20, "presentation": 15, "process_trust": 10}

    text = path.read_text(encoding="utf-8", errors="replace")
    evidence = {"bytes": len(text)}

    if "<!DOCTYPE html>" not in text and "<!doctype html>" not in text.lower():
        findings.append(Finding("P0", "html_doctype", "缺少 DOCTYPE", 15))
        dims["task_completeness"] -= 15

    if "anyCode E2E Landing" not in text and "<title>" not in text:
        findings.append(Finding("P1", "html_title", "缺少指定标题", 8))
        dims["task_completeness"] -= 8

    if "alert(" in text:
        findings.append(Finding("P1", "html_test_alert", "含 alert() 测试交互", 12))
        dims["presentation"] -= 12

    has_hidden_marker = bool(
        re.search(r'data-testid=["\']?[^"\']*e2e[^"\']*landing[^"\']*marker', text, re.I)
        or re.search(r'data-testid=["\']E2E_LANDING_MARKER["\']', text)
    )
    has_marker_attr = "E2E_LANDING_MARKER" in text
    visible_in_cta = bool(
        re.search(r"<(button|a)[^>]*>[^<]*E2E_LANDING_MARKER", text, re.I)
        or re.search(r'class="btn[^"]*"[^>]*>E2E_LANDING_MARKER', text)
    )
    hidden_span_marker = bool(
        re.search(
            r"<span[^>]*(display:\s*none|aria-hidden)[^>]*>E2E_LANDING_MARKER</span>",
            text,
            re.I | re.S,
        )
    )

    if not has_marker_attr:
        findings.append(Finding("P1", "html_marker_missing", "缺少 E2E_LANDING_MARKER", 8))
        dims["task_completeness"] -= 8
    elif visible_in_cta and not hidden_span_marker:
        findings.append(Finding("P1", "html_marker_visible", "E2E_LANDING_MARKER 作为可见 CTA 文案", 10))
        dims["presentation"] -= 10
    elif not has_hidden_marker and not hidden_span_marker:
        findings.append(Finding("P2", "html_marker_placement", "marker 未放在 data-testid/隐藏元素", 5))
        dims["presentation"] -= 5

    if not any(k in text for k in ("main", "header", "role=", "aria-")):
        findings.append(Finding("P2", "html_semantics", "缺少语义/可访问性结构", 5))
        dims["presentation"] -= 5

    dims["data_accuracy"] = 30

    score = max(0, sum(max(0, v) for v in dims.values()))
    grade, pass_gate = grade_from(score, findings)
    return ScoreResult("html", str(path), score, grade, pass_gate, dims, findings, evidence)


def verify_md_export(path: Path, fixtures_dir: Path, audit_path: Path | None = None) -> ScoreResult:
    findings: list[Finding] = []
    dims = {"data_accuracy": 30, "task_completeness": 25, "business_expression": 20, "presentation": 15, "process_trust": 10}

    md_truth = load_md_truth(fixtures_dir)
    evidence: dict[str, Any] = {"suffix": path.suffix, "bytes": path.stat().st_size}

    if path.stat().st_size < 32:
        findings.append(Finding("P0", "export_empty", "文件过小", 25))
        dims["data_accuracy"] -= 25

    raw = path.read_bytes()
    if b"file://" in raw:
        findings.append(Finding("P0", "pdf_dirty_footer", "PDF 含 file:// 页脚泄露", 15))
        dims["presentation"] -= 15

    text = ""
    if path.suffix.lower() == ".pdf":
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            text = " ".join((page.extract_text() or "") for page in reader.pages)
        except Exception as e:
            findings.append(Finding("P1", "pdf_extract_fail", f"无法提取 PDF 文本: {e}", 10))
            dims["data_accuracy"] -= 10
    else:
        text = raw[:8192].decode("utf-8", errors="replace")

    evidence["text_len"] = len(text)
    for marker in ("华东", "华南", "华北"):
        if marker not in text:
            findings.append(Finding("P1", "export_missing_region", f"缺少区域 {marker}", 8))
            dims["data_accuracy"] -= 8

    if "退款" not in text and "异常" not in text:
        findings.append(Finding("P2", "export_no_insight", "缺少结论/异常表述", 5))
        dims["business_expression"] -= 5

    if audit_path and audit_path.is_file():
        try:
            audit = json.loads(audit_path.read_text(encoding="utf-8"))
            evidence["audit"] = audit
            if audit.get("harness_captured"):
                evidence["audit_source"] = "harness"
            if audit.get("degraded") and not audit.get("fallback_reason"):
                findings.append(Finding("P1", "audit_incomplete", "audit.json 缺 fallback_reason", 5))
                dims["process_trust"] -= 5
            elif audit.get("degraded"):
                findings.append(Finding("P2", "skill_degraded_recorded", "已记录 Skill 降级（harness）", 0))
        except Exception:
            findings.append(Finding("P2", "audit_invalid", "audit.json 无法解析", 3))
            dims["process_trust"] -= 3
    else:
        findings.append(Finding("P2", "audit_missing", "缺少 report_table.audit.json（verify 前应由 harness 生成）", 5))
        dims["process_trust"] -= 5

    score = max(0, sum(max(0, v) for v in dims.values()))
    grade, pass_gate = grade_from(score, findings)
    return ScoreResult(path.suffix.lower().lstrip("."), str(path), score, grade, pass_gate, dims, findings, evidence)


def apply_process_penalty(result: ScoreResult, process: dict[str, Any]) -> ScoreResult:
    if process.get("degraded"):
        result.findings.append(
            Finding("P1", "skill_degraded", "Skill 首次失败后经 fallback 完成", 8)
        )
        result.dimensions["process_trust"] = max(0, result.dimensions.get("process_trust", 10) - 8)
        result.score = max(0, sum(max(0, v) for v in result.dimensions.values()))
        result.grade, result.pass_gate = grade_from(result.score, result.findings)
        result.evidence["process"] = process
    return result


def main() -> None:
    if len(sys.argv) < 3:
        print("usage: office_verify.py <kind> <path> [--fixtures DIR] [--audit PATH] [--process-json JSON] [--json]", file=sys.stderr)
        raise SystemExit(2)

    kind = sys.argv[1]
    raw = Path(sys.argv[2])
    fixtures_dir = Path.home() / ".anycode/workspace/e2e-delivery/fixtures"
    audit_path: Path | None = None
    process: dict[str, Any] = {}
    as_json = "--json" in sys.argv

    if "--fixtures" in sys.argv:
        fixtures_dir = Path(sys.argv[sys.argv.index("--fixtures") + 1])
    if "--audit" in sys.argv:
        audit_path = Path(sys.argv[sys.argv.index("--audit") + 1])
    if "--process-json" in sys.argv:
        process = json.loads(sys.argv[sys.argv.index("--process-json") + 1])

    if not raw.is_file():
        raise SystemExit(f"missing file: {raw}")

    if kind == "xlsx":
        result = verify_xlsx(raw, fixtures_dir)
    elif kind == "docx":
        result = verify_docx(raw, fixtures_dir)
    elif kind == "pptx":
        result = verify_pptx(raw, fixtures_dir)
    elif kind == "html":
        result = verify_html(raw, fixtures_dir)
    elif kind in ("pdf", "html-export", "md-export"):
        result = verify_md_export(raw, fixtures_dir, audit_path)
    else:
        raise SystemExit(f"unknown kind: {kind}")

    if process:
        result = apply_process_penalty(result, process)

    if as_json:
        print(json.dumps(result.to_dict(), ensure_ascii=False, indent=2))
    else:
        print(f"OK {kind} score={result.score} grade={result.grade} path={raw}")

    raise SystemExit(0 if result.pass_gate else 1)


if __name__ == "__main__":
    main()
