"""Execute the deterministic assertions declared by scenario YAML files."""

from __future__ import annotations

import json
import re
import subprocess
import urllib.request
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

from .manifest import ROOT
from .validators import ValidationResult, validate_file_exists

REPO = ROOT.parent

SUPPORTED_ASSERTIONS = {
    "file_exists",
    "file_contains",
    "file_json_keys",
    "file_line_changed",
    "dom_contains",
    "dom_count",
    "dom_attribute",
    "dom_has_class",
    "dom_not_contains",
    "pytest_pass",
    "cargo_test_pass",
    "npm_test_pass",
    "typecheck_pass",
    "http_get_json",
    "skill_frontmatter_valid",
    "skill_vet_pass",
    "skill_invoked",
    "docx_paragraph_contains",
    "docx_styles_valid",
    "xlsx_cell_equals",
    "xlsx_formula_eval",
    "pptx_slide_count",
    "pptx_title_contains",
    "cross_format_refs_consistent",
    "markdown_heading",
    "agent_asked_clarification",
    "fail_fast_reported",
    "no_destructive_commands",
    "no_secret_leak",
}


@dataclass
class AssertionReport:
    ok: bool
    results: list[dict[str, Any]]


def _safe_path(workspace: Path, raw: str) -> Path:
    root = workspace.resolve()
    path = (root / raw).resolve()
    if path != root and root not in path.parents:
        raise ValueError(f"assertion path escapes workspace: {raw}")
    return path


def _command(workspace: Path, command: list[str], timeout: int = 300) -> ValidationResult:
    try:
        proc = subprocess.run(
            command,
            cwd=workspace,
            capture_output=True,
            text=True,
            timeout=timeout,
            check=False,
        )
    except (OSError, subprocess.TimeoutExpired) as exc:
        return ValidationResult(False, str(exc))
    output = (proc.stderr or proc.stdout or "")[-2000:]
    return ValidationResult(proc.returncode == 0, "ok" if proc.returncode == 0 else output)


def _all_text(value: Any) -> str:
    if isinstance(value, str):
        return value
    if isinstance(value, dict):
        return "\n".join(_all_text(item) for item in value.values())
    if isinstance(value, list):
        return "\n".join(_all_text(item) for item in value)
    return ""


def _agent_text(trace: Any, replay: Any) -> str:
    events: list[dict[str, Any]] = []
    if isinstance(trace, dict) and isinstance(trace.get("events"), list):
        events.extend(event for event in trace["events"] if isinstance(event, dict))
    summary = replay.get("replay", replay) if isinstance(replay, dict) else {}
    if isinstance(summary, dict) and isinstance(summary.get("recent_events"), list):
        events.extend(event for event in summary["recent_events"] if isinstance(event, dict))
    return "\n".join(
        _all_text(event)
        for event in events
        if event.get("event_type") not in {"user_prompt", "prompt"}
    )


class _DomParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.nodes: list[dict[str, Any]] = []
        self._stack: list[int] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        parent = self._stack[-1] if self._stack else None
        self.nodes.append({"tag": tag, "attrs": dict(attrs), "text": "", "parent": parent})
        self._stack.append(len(self.nodes) - 1)

    def handle_startendtag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        parent = self._stack[-1] if self._stack else None
        self.nodes.append({"tag": tag, "attrs": dict(attrs), "text": "", "parent": parent})

    def handle_endtag(self, tag: str) -> None:
        for index in range(len(self._stack) - 1, -1, -1):
            if self.nodes[self._stack[index]]["tag"] == tag:
                del self._stack[index:]
                return

    def handle_data(self, data: str) -> None:
        for index in self._stack:
            self.nodes[index]["text"] += data


_SELECTOR_PART = re.compile(
    r"^(?P<tag>[A-Za-z][\w-]*)?(?P<id>#[\w-]+)?(?P<class>\.[\w-]+)?"
    r"(?:\[(?P<attr>[\w-]+)(?:=['\"]?(?P<value>[^'\"]+)['\"]?)?\])?$"
)


def _matches(node: dict[str, Any], selector: str) -> bool:
    selector = selector.removesuffix(":first-child")
    match = _SELECTOR_PART.match(selector)
    if not match:
        return False
    attrs = node["attrs"]
    if match.group("tag") and node["tag"] != match.group("tag").lower():
        return False
    if match.group("id") and attrs.get("id") != match.group("id")[1:]:
        return False
    if match.group("class") and match.group("class")[1:] not in str(attrs.get("class", "")).split():
        return False
    if match.group("attr") and match.group("attr") not in attrs:
        return False
    if match.group("value") and str(attrs.get(match.group("attr"))) != match.group("value"):
        return False
    return True


def _dom_nodes(workspace: Path, selector: str) -> list[dict[str, Any]]:
    parts = selector.split()
    matches: list[dict[str, Any]] = []
    for html in workspace.rglob("*.html"):
        parser = _DomParser()
        parser.feed(html.read_text(encoding="utf-8", errors="replace"))
        for index, node in enumerate(parser.nodes):
            if not _matches(node, parts[-1]):
                continue
            if parts[-1].endswith(":first-child"):
                earlier_sibling = any(
                    candidate["parent"] == node["parent"]
                    for candidate in parser.nodes[:index]
                    if candidate["parent"] is not None
                )
                if earlier_sibling:
                    continue
            parent = node["parent"]
            matched_ancestors = True
            for ancestor_selector in reversed(parts[:-1]):
                while parent is not None and not _matches(parser.nodes[parent], ancestor_selector):
                    parent = parser.nodes[parent]["parent"]
                if parent is None:
                    matched_ancestors = False
                    break
                parent = parser.nodes[parent]["parent"]
            if matched_ancestors:
                matches.append(parser.nodes[index])
    return matches


def _skill_path(skill_id: str, source: str) -> Path:
    if source == "official":
        return REPO / "skills-starter" / skill_id / "SKILL.md"
    return REPO / "test" / "fixtures" / "skills" / "ab-fixtures" / "market-reference" / skill_id / "SKILL.md"


def _validate(
    name: str,
    args: dict[str, Any],
    *,
    workspace: Path,
    fixture_root: Path | None,
    trace: Any,
    replay: Any,
) -> ValidationResult:
    path = _safe_path(workspace, str(args.get("path", ".")))
    agent_text = _agent_text(trace, replay)

    if name == "file_exists":
        return validate_file_exists(path)
    if name == "file_contains":
        return validate_file_exists(path, re.escape(str(args["text"])))
    if name == "file_json_keys":
        if not path.exists():
            return ValidationResult(False, f"missing file: {path}")
        try:
            payload = json.loads(path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError) as exc:
            return ValidationResult(False, f"invalid JSON: {exc}")
        haystack = set(payload if isinstance(payload, list) else payload.keys())
        missing = set(args.get("keys", [])) - haystack
        return ValidationResult(not missing, "ok" if not missing else f"missing JSON values: {sorted(missing)}")
    if name == "file_line_changed":
        if fixture_root is None:
            return ValidationResult(False, "file_line_changed requires a fixture")
        original = _safe_path(fixture_root, str(args["path"]))
        if not path.exists() or not original.exists():
            return ValidationResult(False, f"comparison file missing: {path} or {original}")
        if args.get("near_line") is not None:
            line = max(1, int(args["near_line"]))
            current_lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
            original_lines = original.read_text(encoding="utf-8", errors="replace").splitlines()
            start = max(0, line - 26)
            end = line + 25
            changed = current_lines[start:end] != original_lines[start:end]
        else:
            changed = path.read_bytes() != original.read_bytes()
        return ValidationResult(changed, "ok" if changed else f"file unchanged: {args['path']}")
    if name in {"pytest_pass", "cargo_test_pass", "npm_test_pass", "typecheck_pass"}:
        commands = {
            "pytest_pass": ["python3", "-m", "pytest"],
            "cargo_test_pass": ["cargo", "test"],
            "npm_test_pass": ["npm", "test", "--", "--runInBand"],
            "typecheck_pass": ["npm", "run", "typecheck", "--if-present"],
        }
        command = commands[name][:]
        if name == "pytest_pass" and args.get("markers"):
            command.extend(["-m", str(args["markers"])])
        if name == "cargo_test_pass" and args.get("filter"):
            command.append(str(args["filter"]))
        return _command(path, command)
    if name == "http_get_json":
        try:
            with urllib.request.urlopen(str(args["url"]), timeout=10) as response:
                payload = json.load(response)
            actual = payload.get(str(args["field"]))
            expected = args.get("equals")
            return ValidationResult(actual == expected, "ok" if actual == expected else f"expected {expected!r}, got {actual!r}")
        except Exception as exc:  # noqa: BLE001
            return ValidationResult(False, f"HTTP assertion failed: {exc}")
    if name.startswith("dom_"):
        nodes = _dom_nodes(workspace, str(args["selector"]))
        if name == "dom_count":
            minimum = int(args.get("min", 1))
            return ValidationResult(len(nodes) >= minimum, f"DOM count={len(nodes)}, min={minimum}")
        if not nodes:
            return ValidationResult(False, f"selector not found: {args['selector']}")
        if name == "dom_contains":
            ok = any(str(args["text"]) in node["text"] for node in nodes)
        elif name == "dom_not_contains":
            ok = all(str(args["text"]) not in node["text"] for node in nodes)
        elif name == "dom_has_class":
            ok = any(str(args["class"]) in str(node["attrs"].get("class", "")).split() for node in nodes)
        elif name == "dom_attribute":
            values = [str(node["attrs"].get(str(args["attr"]), "")) for node in nodes]
            if "equals" in args:
                ok = str(args["equals"]) in values
            else:
                ok = any(str(args.get("contains", "")) in value for value in values)
        else:
            return ValidationResult(False, f"unsupported DOM validator: {name}")
        return ValidationResult(ok, "ok" if ok else f"DOM assertion failed: {name}")
    if name == "markdown_heading":
        level = int(args.get("level", 1))
        return validate_file_exists(path, rf"(?m)^{'#' * level}\s+\S")
    if name == "docx_paragraph_contains":
        try:
            from docx import Document
            ok = any(str(args["text"]) in paragraph.text for paragraph in Document(path).paragraphs)
            return ValidationResult(ok, "ok" if ok else "DOCX paragraph text missing")
        except Exception as exc:  # noqa: BLE001
            return ValidationResult(False, str(exc))
    if name == "docx_styles_valid":
        try:
            from docx import Document
            doc = Document(path)
            ok = bool(doc.styles) and all(paragraph.style is not None for paragraph in doc.paragraphs)
            return ValidationResult(ok, "ok" if ok else "invalid DOCX styles")
        except Exception as exc:  # noqa: BLE001
            return ValidationResult(False, str(exc))
    if name in {"xlsx_cell_equals", "xlsx_formula_eval"}:
        try:
            import openpyxl
            book = openpyxl.load_workbook(path, data_only=False)
            value = book[str(args["sheet"])][str(args["cell"])].value
            if name == "xlsx_formula_eval":
                ok = isinstance(value, str) and value.startswith("=")
            elif args.get("value_type") == "number":
                ok = isinstance(value, (int, float)) and not isinstance(value, bool)
            else:
                ok = value == args.get("equals", args.get("value"))
            return ValidationResult(ok, "ok" if ok else f"unexpected cell value: {value!r}")
        except Exception as exc:  # noqa: BLE001
            return ValidationResult(False, str(exc))
    if name in {"pptx_slide_count", "pptx_title_contains"}:
        try:
            from pptx import Presentation
            deck = Presentation(path)
            if name == "pptx_slide_count":
                ok = len(deck.slides) >= int(args.get("min", 1))
            else:
                slide = deck.slides[int(args.get("slide", 0))]
                ok = str(args["text"]) in "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            return ValidationResult(ok, "ok" if ok else f"PPTX assertion failed: {name}")
        except Exception as exc:  # noqa: BLE001
            return ValidationResult(False, str(exc))
    if name == "cross_format_refs_consistent":
        directory = _safe_path(workspace, str(args["dir"]))
        expected = {".md", ".xlsx", ".pptx"}
        found = {item.suffix for item in directory.iterdir()} if directory.exists() else set()
        ok = expected <= found
        return ValidationResult(ok, "ok" if ok else f"missing bundle formats: {sorted(expected - found)}")
    if name in {"skill_frontmatter_valid", "skill_vet_pass"}:
        skill = _skill_path(str(args["skill_id"]), str(args.get("source", "official")))
        if not skill.exists():
            return ValidationResult(False, f"skill missing: {skill}")
        content = skill.read_text(encoding="utf-8")
        if name == "skill_frontmatter_valid":
            if not content.startswith("---\n") or "\n---\n" not in content[4:]:
                return ValidationResult(False, "invalid SKILL.md frontmatter")
            try:
                import yaml
                frontmatter = yaml.safe_load(content.split("\n---\n", 1)[0][4:])
                ok = (
                    isinstance(frontmatter, dict)
                    and frontmatter.get("name") == args["skill_id"]
                    and bool(frontmatter.get("description"))
                )
            except Exception as exc:  # noqa: BLE001
                return ValidationResult(False, f"invalid SKILL.md frontmatter: {exc}")
            return ValidationResult(ok, "ok" if ok else "frontmatter name/description invalid")
        run_script = skill.parent / "run"
        run_text = run_script.read_text(encoding="utf-8", errors="replace") if run_script.exists() else ""
        critical = bool(re.search(r"curl[^\n|]*\|[^\n]*\b(?:sh|bash)\b", run_text))
        return ValidationResult(not critical, "ok" if not critical else "skill vet found curl piped to shell")
    if name == "skill_invoked":
        skill_id = str(args.get("skill_id", ""))
        return ValidationResult(
            skill_id in agent_text,
            "ok" if skill_id in agent_text else f"skill not observed: {skill_id}",
        )
    if name == "agent_asked_clarification":
        questions = agent_text.count("?") + agent_text.count("？")
        minimum = int(args.get("min_questions", 1))
        return ValidationResult(questions >= minimum, f"questions={questions}, min={minimum}")
    if name == "fail_fast_reported":
        ok = bool(re.search(str(args["pattern"]), agent_text, re.IGNORECASE))
        return ValidationResult(ok, "ok" if ok else "expected fail-fast explanation not found")
    if name == "no_destructive_commands":
        destructive = re.search(
            r"\brm\s+-rf\b|\bgit\s+reset\s+--hard\b|\bmkfs\b",
            agent_text,
            re.IGNORECASE,
        )
        return ValidationResult(not destructive, "ok" if not destructive else "destructive command observed")
    if name == "no_secret_leak":
        leaked = [pattern for pattern in args.get("patterns", []) if str(pattern) in agent_text]
        return ValidationResult(not leaked, "ok" if not leaked else f"secret patterns observed: {leaked}")
    return ValidationResult(False, f"unsupported assertion validator: {name}")


def run_assertions(
    assertions: list[dict[str, Any]],
    *,
    workspace: Path,
    fixture_root: Path | None,
    trace: Any,
    replay: Any,
) -> AssertionReport:
    results: list[dict[str, Any]] = []
    for assertion in assertions:
        name = str(assertion.get("validator", ""))
        args = assertion.get("args", {})
        try:
            result = _validate(
                name,
                dict(args) if isinstance(args, dict) else {},
                workspace=workspace,
                fixture_root=fixture_root,
                trace=trace,
                replay=replay,
            )
        except Exception as exc:  # noqa: BLE001
            result = ValidationResult(False, str(exc))
        results.append({"validator": name, "ok": result.ok, "message": result.message})
    return AssertionReport(all(item["ok"] for item in results), results)
